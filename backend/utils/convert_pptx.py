import os
import subprocess
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
import textwrap

UPLOADS_DIR = "/tmp/slideconfirm_uploads"
Path(UPLOADS_DIR).mkdir(parents=True, exist_ok=True)

def convert_pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """Конвертирует PPTX в изображения JPG"""
    
    # Первый приоритет - LibreOffice
    libreoffice_path = find_libreoffice()
    if libreoffice_path:
        try:
            return convert_with_libreoffice(libreoffice_path, pptx_path, output_dir)
        except Exception as e:
            print(f"LibreOffice конвертирование не удалось: {e}")
            # Fallback на python-pptx
    
    # Fallback на python-pptx с улучшенным отображением
    return convert_with_enhanced_pptx(pptx_path, output_dir)


def find_libreoffice():
    """Ищет путь к LibreOffice"""
    possible_paths = [
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        '/usr/bin/libreoffice',
        '/usr/bin/soffice',
        '/opt/libreoffice/program/soffice',
        '/usr/local/bin/libreoffice',
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            return path
    
    # Пытаемся найти через which
    try:
        result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True, timeout=5)
        if result.returncode == 0:
            return result.stdout.strip()
    except:
        pass
    
    return None


def convert_with_libreoffice(libreoffice_path: str, pptx_path: str, output_dir: str) -> list[str]:
    """Конвертирует PPTX в PDF через LibreOffice, затем PDF в JPG"""
    try:
        pdf_path = os.path.join(output_dir, "presentation.pdf")
        
        # Конвертируем PPTX в PDF через LibreOffice
        subprocess.run([
            libreoffice_path, '--headless', '--convert-to', 'pdf',
            '--outdir', output_dir, pptx_path
        ], check=True, capture_output=True, timeout=300)
        
        # Конвертируем PDF в JPG через pdftoppm
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path, dpi=150, fmt='jpg')
        
        # Сохраняем изображения
        image_paths = []
        for idx, image in enumerate(images, 1):
            image_path = os.path.join(output_dir, f"slide_{idx:03d}.jpg")
            image.save(image_path, 'JPEG', quality=90)
            image_paths.append(image_path)
        
        # Удаляем PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        
        return image_paths
        
    except Exception as e:
        raise Exception(f"LibreOffice conversion failed: {str(e)}")


def convert_with_enhanced_pptx(pptx_path: str, output_dir: str) -> list[str]:
    """Улучшенное конвертирование с использованием python-pptx"""
    try:
        from pptx import Presentation
    except ImportError:
        raise Exception("python-pptx не установлена")
    
    try:
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            raise Exception("PPTX файл не содержит слайдов")
        
        image_paths = []
        width, height = 1280, 720
        
        # Для каждого слайда создаем JPG
        for idx, slide in enumerate(prs.slides, 1):
            # Создаем белое изображение
            img = Image.new('RGB', (width, height), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Рисуем серый фон для разделения
            draw.rectangle([0, 0, width, height], fill=(250, 250, 250), outline=(150, 150, 200), width=2)
            
            # Добавляем содержимое слайда
            y_offset = 30
            shape_count = 0
            
            # Обрабатываем все shape'ы на слайде
            for shape_idx, shape in enumerate(slide.shapes):
                # Обработка текстовых блоков
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    if text_frame.text.strip():
                        # Анализируем структуру текста
                        for para_idx, paragraph in enumerate(text_frame.paragraphs):
                            text = paragraph.text.strip()
                            if not text:
                                continue
                            
                            level = paragraph.level
                            
                            # Определяем стиль в зависимости от уровня
                            if level == 0 and para_idx == 0:
                                # Заголовок слайда
                                font_size = 18
                                color = (0, 0, 100)
                                y_offset += 10
                            elif level == 0:
                                # Основной текст
                                font_size = 14
                                color = (30, 30, 30)
                            else:
                                # Текст второго уровня
                                font_size = 12
                                color = (80, 80, 80)
                                text = "  • " + text
                            
                            # Рисуем текст с переносами
                            lines = textwrap.wrap(text, width=90)
                            for line in lines:
                                if y_offset > height - 80:
                                    break
                                
                                draw.text((40, y_offset), line, fill=color, font=None)
                                y_offset += font_size + 4
                            
                            shape_count += 1
                            
                            if y_offset > height - 80:
                                break
                
                # Обработка таблиц
                elif hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    y_offset += 15
                    
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip()
                            if text:
                                color = (0, 0, 100) if row_idx == 0 else (50, 50, 50)
                                x_pos = 40 + col_idx * 280
                                
                                # Обрезаем длинный текст
                                if len(text) > 30:
                                    text = text[:27] + "..."
                                
                                draw.text((x_pos, y_offset), text, fill=color, font=None)
                        
                        y_offset += 25
                        if y_offset > height - 80:
                            break
                    
                    shape_count += 1
                
                if y_offset > height - 80:
                    break
            
            # Если слайд пустой, добавляем сообщение
            if shape_count == 0:
                draw.text((width // 2 - 100, height // 2 - 30), 
                         f"Слайд {idx} (без текста)", 
                         fill=(150, 150, 150), font=None)
            
            # Добавляем номер слайда
            draw.text((width - 150, height - 40), f"Слайд {idx}", 
                     fill=(150, 150, 150), font=None)
            
            # Сохраняем JPG
            image_path = os.path.join(output_dir, f"slide_{idx:03d}.jpg")
            img.save(image_path, 'JPEG', quality=95)
            image_paths.append(image_path)
        
        return image_paths
        
    except Exception as e:
        raise Exception(f"PPTX conversion failed: {str(e)}")


def draw_wrapped_text(draw, text, x, y, max_width, color, font_size=14, bold=False):
    """Рисует текст с переносом строк"""
    # Подготовляем текст
    lines = text.split('\n')
    
    try:
        if font_size > 16:
            font = ImageFont.load_default()
        else:
            font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    current_y = y
    char_width = 8  # Примерная ширина символа
    chars_per_line = max(1, max_width // char_width)
    
    for line in lines:
        if not line.strip():
            current_y += font_size + 5
            continue
        
        # Разбиваем длинные строки
        wrapped_lines = textwrap.wrap(line.strip(), width=chars_per_line)
        
        for wrapped_line in wrapped_lines:
            if current_y > 700:  # Не выходим за границы
                return current_y
            
            draw.text((x, current_y), wrapped_line, fill=color, font=font)
            current_y += font_size + 5
    
    return current_y
